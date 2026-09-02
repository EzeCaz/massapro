"""Rebrand ConnexAI Enterprise Product Brochure (2026).pptx as a MassaPro deck.

Approach
--------
1. Load the source PPTX with python-pptx.
2. Walk every slide + every shape (recursing into GROUP shapes) and:
   a. Replace text in every paragraph using a brand-text mapping (ConnexAI → MassaPro, etc.)
      while preserving the original run formatting (font, size, color, bold).
   b. Recolor text runs from ConnexAI teal/green → MassaPro purple.
3. Recolor theme solidFills: teal family → MassaPro purple family.
4. Save as new PPTX.

This preserves all original layout, images, and product screenshots — only the
text-level branding and the brand color tokens are swapped.
"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import copy
import re

SRC = '/home/z/my-project/upload/(US) ConnexAI Enterprise Product Brochure (2026).pptx'
DST = '/home/z/my-project/download/Massapro-Enterprise-Brochure-Rebranded.pptx'

# ---------- Brand text replacements ----------
# Apply per-paragraph (full-paragraph match) AND per-run (substring match).
# Per-paragraph is safer for headlines that are split into multiple runs.
PARAGRAPH_REPLACEMENTS = {
    # Cover slide
    'ConnexAl': 'MassaPro',          # the "I" is rendered as lowercase L in the original
    'ConnexAI': 'MassaPro',
    'Conversational Agentic AI': 'Agentic AI Customer Journey Platform',
    'The Global Leader in Conversational': 'The Global Leader in Agentic AI',
    # Closing slide
    'www.connexai': 'www.massapro.com',
    'hello@connexai': 'hello@massapro.com',
    # Section headers
    'ConnexAI is the Leading AI Model for Voice Cloning': 'MassaPro is the Leading AI Model for Voice Cloning',
    "ConnexAI's cutting-edge Text-to-Speech": "MassaPro's cutting-edge Text-to-Speech",
    # Generic footer / brand mentions
    'ConnexAI Conversations': 'MassaPro Conversations',
    'ConnexAI AI Conversations': 'MassaPro AI Conversations',
    'ConnexAI AI Agent': 'MassaPro AI Agent',
    'ConnexAI dashboards': 'MassaPro dashboards',
    "ConnexAI's drag-and-drop": "MassaPro's drag-and-drop",
    'ConnexAI suite': 'MassaPro suite',
    'connect ConnexAI to': 'connect MassaPro to',
    'platform handles the entire interaction': 'platform orchestrates the entire customer journey',
    'ground every agent in a knowledge base you fully own and control': 'ground every agent in a knowledge base you fully own and control, orchestrated by our proprietary LLM',
    # Slide 2 overview — pivot the messaging toward MassaPro's value prop
    'Unlike most providers that rely on licensed, general-purpose models, we design and control every layer of our AI technology.':
        'Unlike most providers that charge per-minute AI tokens for voice or text, MassaPro offers a flat pay-per-agent + volume pricing model — scaling without surprise bills.',
    'We design and own every layer of the AI Stack, ensuring best in class performance and data sovereignty for the enterprise customers using our technology.':
        'Our proprietary LLM orchestrates the entire customer journey across every channel — replacing the call center with a growth machine powered by Agentic AI.',
    'One Platform. A Complete AI Workforce.': 'One Platform. A Complete Agentic AI Workforce.',
}

# Substring replacements applied per-run (only if the paragraph replacement didn't match)
RUN_SUBSTRING_REPLACEMENTS = [
    ('ConnexAI', 'MassaPro'),
    ('ConnexAl', 'MassaPro'),  # stylized lowercase-L
    ('connexai', 'massapro'),
    ('ConnexAI', 'MassaPro'),
    ('Athena LLM', 'MassaPro LLM'),
    ('Athena\'s Large Language Model', 'Our proprietary Large Language Model'),
    ('Athena uses reasoning', 'Our proprietary LLM uses reasoning'),
    ('Athena', 'MassaPro'),   # last-resort Athena → MassaPro
]

# ---------- Color replacements ----------
# ConnexAI teal/green family → MassaPro purple family
# Original brand colors observed (from slide XML):
#   #49B395, #5DA391, #3F796B, #249A87, #369383, #427062, #6B8E87 — teal/green
#   Dark variants: #000E0F, #011111, #031618 — near-black teal
# MassaPro purple family:
#   Primary: #7C3AED  Dark: #6D28D9  Darker: #5B21B6
#   Light: #A78BFA    Soft: #C4B5FD  100: #EDE9FE
#   Bg dark: #1A0B2E  Bg darker: #0F0518
COLOR_MAP = {
    # Bright brand teal
    '49B395': '7C3AED',
    '5DA391': '7C3AED',
    '3F796B': '6D28D9',
    '249A87': '7C3AED',
    '369383': '6D28D9',
    '427062': '5B21B6',
    '6B8E87': 'A78BFA',
    # Near-black teal (background tints) → near-black purple
    '000E0F': '0F0518',
    '011111': '1A0B2E',
    '031618': '1A0B2E',
    '001F1D': '1A0B2E',
    '00241F': '1A0B2E',
    # Accent green tones if any
    'D8FFFF': 'EDE9FE',
    'E2FBFB': 'F5F3FF',
}

def normalize_color(hex_str):
    """Normalize a hex color string to uppercase 6-char without #."""
    if not hex_str:
        return None
    h = hex_str.upper().lstrip('#')
    if len(h) == 6 and all(c in '0123456789ABCDEF' for c in h):
        return h
    return None

def map_color(hex_str):
    """Map original brand colors to MassaPro equivalents. Returns None if no change."""
    h = normalize_color(hex_str)
    if h and h in COLOR_MAP:
        return COLOR_MAP[h]
    return None

# ---------- Helpers ----------
def iter_text_frames(shapes):
    for s in shapes:
        if s.shape_type == 6:  # GROUP
            yield from iter_text_frames(s.shapes)
        elif s.has_text_frame:
            yield s, s.text_frame

def _norm(s):
    return s.replace('\x0b', '').replace('\r', '').strip()

def replace_in_paragraph_preserve_format(paragraph, new_text):
    """Replace paragraph text using the first run's formatting."""
    runs = paragraph.runs
    if not runs:
        paragraph.add_run().text = new_text
        return
    runs[0].text = new_text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)

def apply_paragraph_replacements(text_frame):
    """For each paragraph, try to match the full normalized text against
    PARAGRAPH_REPLACEMENTS. If matched, replace preserving formatting."""
    m = {_norm(k): v for k, v in PARAGRAPH_REPLACEMENTS.items()}
    for p in text_frame.paragraphs:
        key = _norm(p.text)
        if key and key in m:
            replace_in_paragraph_preserve_format(p, m[key])

def apply_run_substring_replacements(text_frame):
    """For each run, apply RUN_SUBSTRING_REPLACEMENTS (substring search/replace).
    Preserves all run formatting."""
    for p in text_frame.paragraphs:
        for r in p.runs:
            original = r.text
            new = original
            for old, repl in RUN_SUBSTRING_REPLACEMENTS:
                if old in new:
                    new = new.replace(old, repl)
            if new != original:
                r.text = new

def recolor_runs(text_frame):
    """For each run, if its color is in COLOR_MAP, swap to MassaPro equivalent."""
    for p in text_frame.paragraphs:
        for r in p.runs:
            try:
                if r.font.color and r.font.color.type is not None:
                    rgb = r.font.color.rgb
                    if rgb:
                        hex_str = str(rgb).upper().lstrip('#')
                        if hex_str in COLOR_MAP:
                            r.font.color.rgb = RGBColor.from_string(COLOR_MAP[hex_str])
            except Exception:
                # Color may be theme-based or inherited — skip
                pass

def recolor_shape_fills(shape):
    """Walk shape XML and recolor any <a:srgbClr> with values in COLOR_MAP."""
    try:
        spPr = shape._element.spPr
        if spPr is None:
            return
        for srgbClr in spPr.iter(qn('a:srgbClr')):
            val = srgbClr.get('val', '').upper()
            if val in COLOR_MAP:
                srgbClr.set('val', COLOR_MAP[val])
    except Exception:
        pass

def walk_and_rebrand(shapes):
    """Recursively walk shapes, applying text replacements + color swaps."""
    for s in shapes:
        if s.shape_type == 6:  # GROUP
            walk_and_rebrand(s.shapes)
            continue
        if s.has_text_frame:
            apply_paragraph_replacements(s.text_frame)
            apply_run_substring_replacements(s.text_frame)
            recolor_runs(s.text_frame)
        recolor_shape_fills(s)

# ---------- Main ----------
prs = Presentation(SRC)
print(f'Loaded {SRC}')
print(f'Slides: {len(prs.slides)}')

for i in range(len(prs.slides._sldIdLst)):
    slide = prs.slides[i]
    walk_and_rebrand(slide.shapes)
    print(f'  ✓ Slide {i+1} rebranded')

prs.save(DST)
print(f'\nSaved: {DST}')
