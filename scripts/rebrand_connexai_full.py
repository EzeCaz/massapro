"""Comprehensive MassaPro rebrand of ConnexAI Enterprise Product Brochure.

Approach (mirrors the IVR Telepresencia rebrand that worked well):
1. Load source PPTX with python-pptx
2. For every text run, apply brand text replacements AND color swaps:
   * "ConnexAI" / "ConnexAl" → "MassaPro"
   * "Athena" → "MassaPro" (LLM references)
   * "Athena LLM" → "MassaPro LLM"
   * "Interactive Powers" → "MassaPro" (if any)
   * "www.connexai" / "connexai.com" → "massapro.com"
   * Email "hello@connexai" → "hello@massapro.com"
   * Slide 2 cover text: pivot messaging to MassaPro value-prop
3. Recolor all ConnexAI teal/green → MassaPro purple:
   * #49B395, #5DA391, #3F796B, #249A87, #369383, #427062, #6B8E87 → #7C3AED
   * Light tints #D8FFFF, #E2FBFB → #EDE9FE / #F5F3FF
4. Recolor dark teal backgrounds → MassaPro dark purple #1A0B2E
5. Add MassaPro brand bar (purple pill, top-right corner) on every body slide
6. Save as new PPTX + Compact (image-optimized) versions
"""
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy
import os

SRC = '/home/z/my-project/upload/(US) ConnexAI Enterprise Product Brochure (2026).pptx'
DST = '/home/z/my-project/download/connexai_massapro_full.pptx'

# ---------- Brand text replacements (paragraph-level, full match) ----------
PARAGRAPH_REPLACEMENTS = {
    # Cover slide
    'ConnexAl': 'MassaPro',
    'ConnexAI': 'MassaPro',
    'Conversational Agentic AI': 'Agentic AI Customer Journey Platform',
    'The Global Leader in Conversational': 'The Global Leader in Agentic AI',
    # Slide 2 overview — pivot messaging to MassaPro value-prop
    'One Platform. A Complete AI Workforce.':
        'One Platform. A Complete Agentic AI Workforce.',
    'Unlike most providers that rely on licensed, general-purpose models, we design and control every layer of our AI technology.':
        'Unlike most providers that charge per-minute AI tokens for voice or text, MassaPro offers a flat pay-per-agent + volume pricing model — scaling without surprise bills.',
    'We design and own every layer of the AI Stack, ensuring best in class performance and data sovereignty for the enterprise customers using our technology.':
        'Our proprietary LLM orchestrates the entire customer journey across every channel — replacing the call center with a growth machine powered by Agentic AI.',
    # Athena → MassaPro LLM
    'Athena LLM': 'MassaPro LLM',
    "Athena's Large Language Model": "Our proprietary Large Language Model",
    'Athena uses reasoning': 'Our proprietary LLM uses reasoning',
    # Closing slide
    'www.connexai': 'www.massapro.com',
    'hello@connexai': 'hello@massapro.com',
    'ConnexAI is the Leading AI Model for Voice Cloning':
        'MassaPro is the Leading AI Model for Voice Cloning',
    "ConnexAI's cutting-edge Text-to-Speech":
        "MassaPro's cutting-edge Text-to-Speech",
    'ConnexAI Conversations': 'MassaPro Conversations',
    'ConnexAI AI Conversations': 'MassaPro AI Conversations',
    'ConnexAI AI Agent': 'MassaPro AI Agent',
    'ConnexAI dashboards': 'MassaPro dashboards',
    "ConnexAI's drag-and-drop": "MassaPro's drag-and-drop",
    'ConnexAI suite': 'MassaPro suite',
    'connect ConnexAI to': 'connect MassaPro to',
    'platform handles the entire interaction':
        'platform orchestrates the entire customer journey',
    'ground every agent in a knowledge base you fully own and control':
        'ground every agent in a knowledge base you fully own and control, orchestrated by our proprietary LLM',
}

# ---------- Substring replacements (per-run) ----------
RUN_SUBSTRING_REPLACEMENTS = [
    ('ConnexAI', 'MassaPro'),
    ('ConnexAl', 'MassaPro'),
    ('connexai', 'massapro'),
    ('Athena LLM', 'MassaPro LLM'),
    ('Athena', 'MassaPro'),
    ('Interactive Powers', 'MassaPro'),
]

# ---------- Color map: ConnexAI teal → MassaPro purple ----------
COLOR_MAP = {
    # Bright brand teal → MassaPro primary purple
    '49B395': '7C3AED',
    '5DA391': '7C3AED',
    '3F796B': '6D28D9',
    '249A87': '7C3AED',
    '369383': '6D28D9',
    '427062': '5B21B6',
    '6B8E87': 'A78BFA',
    # Near-black teal (background tints) → MassaPro dark purple
    '000E0F': '0F0518',
    '011111': '1A0B2E',
    '031618': '1A0B2E',
    '001F1D': '1A0B2E',
    '00241F': '1A0B2E',
    # Light cyan tints → light purple
    'D8FFFF': 'EDE9FE',
    'E2FBFB': 'F5F3FF',
}

def normalize_color(hex_str):
    if not hex_str: return None
    h = hex_str.upper().lstrip('#')
    if len(h) == 6 and all(c in '0123456789ABCDEF' for c in h):
        return h
    return None

def map_color(hex_str):
    h = normalize_color(hex_str)
    if h and h in COLOR_MAP:
        return COLOR_MAP[h]
    return None

# ---------- Helpers ----------
def iter_text_frames(shapes):
    for s in shapes:
        if s.shape_type == 6:  # GROUP
            yield from iter_text_frames(s.shapes)
        elif s.has_text_frame:
            yield s, s.text_frame

def _norm(s):
    return s.replace('\x0b', '').replace('\r', '').strip()

def replace_in_paragraph_preserve_format(paragraph, new_text):
    runs = paragraph.runs
    if not runs:
        paragraph.add_run().text = new_text
        return
    runs[0].text = new_text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)

def apply_paragraph_replacements(text_frame):
    m = {_norm(k): v for k, v in PARAGRAPH_REPLACEMENTS.items()}
    for p in text_frame.paragraphs:
        key = _norm(p.text)
        if key and key in m:
            replace_in_paragraph_preserve_format(p, m[key])

def apply_run_substring_replacements(text_frame):
    for p in text_frame.paragraphs:
        for r in p.runs:
            original = r.text
            new = original
            for old, repl in RUN_SUBSTRING_REPLACEMENTS:
                if old in new:
                    new = new.replace(old, repl)
            if new != original:
                r.text = new

def recolor_runs(text_frame):
    for p in text_frame.paragraphs:
        for r in p.runs:
            try:
                if r.font.color and r.font.color.type is not None:
                    rgb = r.font.color.rgb
                    if rgb:
                        hex_str = str(rgb).upper().lstrip('#')
                        if hex_str in COLOR_MAP:
                            r.font.color.rgb = RGBColor.from_string(COLOR_MAP[hex_str])
            except Exception:
                pass

def recolor_shape_fills(shape):
    """Walk shape XML and recolor any <a:srgbClr> with values in COLOR_MAP."""
    try:
        spPr = shape._element.spPr
        if spPr is None: return
        for srgbClr in spPr.iter(qn('a:srgbClr')):
            val = srgbClr.get('val', '').upper()
            if val in COLOR_MAP:
                srgbClr.set('val', COLOR_MAP[val])
    except Exception:
        pass

def walk_and_rebrand(shapes):
    for s in shapes:
        if s.shape_type == 6:  # GROUP
            walk_and_rebrand(s.shapes)
            continue
        if s.has_text_frame:
            apply_paragraph_replacements(s.text_frame)
            apply_run_substring_replacements(s.text_frame)
            recolor_runs(s.text_frame)
        recolor_shape_fills(s)

# ---------- MassaPro brand bar overlay ----------
def add_massapro_brand_bar(slide, slide_w, slide_h, position='top-right'):
    """Add a small MassaPro pill badge to the slide corner."""
    mp_purple = RGBColor(0x7C, 0x3A, 0xED)
    mp_white  = RGBColor(0xFF, 0xFF, 0xFF)

    badge_w = Inches(1.55)
    badge_h = Inches(0.32)
    if position == 'top-right':
        badge_x = slide_w - badge_w - Inches(0.15)
        badge_y = Inches(0.12)
    elif position == 'top-left':
        badge_x = Inches(0.15)
        badge_y = Inches(0.12)
    else:
        badge_x = slide_w - badge_w - Inches(0.15)
        badge_y = Inches(0.12)

    # Create rounded rectangle badge with white fill + purple border
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, badge_x, badge_y, badge_w, badge_h)
    badge.fill.solid()
    badge.fill.fore_color.rgb = mp_white
    badge.line.color.rgb = mp_purple
    badge.line.width = Pt(1.0)

    # Remove shadow
    spPr = badge._element.spPr
    for tag in ['a:effectLst']:
        existing = spPr.find(qn(tag))
        if existing is not None:
            spPr.remove(existing)
    effectLst = spPr.makeelement(qn('a:effectLst'), {})
    spPr.append(effectLst)

    tf = badge.text_frame
    tf.text = 'MassaPro'
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = 'Inter'
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = mp_purple
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    # Add small purple icon with "M" before text
    icon_size = Inches(0.18)
    icon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        badge_x + Inches(0.08), badge_y + (badge_h - icon_size) // 2,
        icon_size, icon_size)
    icon.fill.solid()
    icon.fill.fore_color.rgb = mp_purple
    icon.line.fill.background()
    icon_tf = icon.text_frame
    icon_tf.text = 'M'
    icon_p = icon_tf.paragraphs[0]
    icon_p.alignment = PP_ALIGN.CENTER
    icon_run = icon_p.runs[0]
    icon_run.font.size = Pt(10)
    icon_run.font.bold = True
    icon_run.font.color.rgb = mp_white
    icon_run.font.name = 'Inter'

# ---------- Main ----------
prs = Presentation(SRC)
print(f'Loaded {SRC}')
print(f'Slides: {len(prs.slides)}')

slide_w_emu = prs.slide_width
slide_h_emu = prs.slide_height

for i in range(len(prs.slides._sldIdLst)):
    slide = prs.slides[i]
    walk_and_rebrand(slide.shapes)
    # Add MassaPro brand bar to every body slide (skip cover slide 1)
    if i > 0:
        add_massapro_brand_bar(slide, slide_w_emu, slide_h_emu, 'top-right')
    print(f'  ✓ Slide {i+1} rebranded')

prs.save(DST)
print(f'\nSaved: {DST}')
print(f'Size: {os.path.getsize(DST)//1024//1024} MB')
